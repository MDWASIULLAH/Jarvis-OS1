"""
capabilities/document_writer.py

Implements JARVIS Section 2.8 (write side) -- generating Word, PowerPoint,
and Excel files from structured content.
"""

from __future__ import annotations

# Office file generation is optional. Previously these were hard imports, so a
# missing python-pptx crashed the whole server at startup rather than just
# disabling .pptx export. Each now fails only when that format is requested.
try:
    from docx import Document
except ImportError:  # pragma: no cover
    Document = None

try:
    from openpyxl import Workbook
except ImportError:  # pragma: no cover
    Workbook = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None


def _require(obj, package: str, fmt: str):
    if obj is None:
        raise RuntimeError(
            f"Writing {fmt} files needs the '{package}' package. "
            f"Install it with: pip install {package}"
        )


def write_docx(path: str, title: str, paragraphs: list) -> str:
    _require(Document, "python-docx", ".docx")
    doc = Document()
    doc.add_heading(title, level=1)
    for p in paragraphs:
        doc.add_paragraph(p)
    doc.save(path)
    return path


def write_pptx(path: str, title: str, slides: list) -> str:
    """slides: list of {'heading': str, 'bullets': [str, ...]}"""
    _require(Presentation, "python-pptx", ".pptx")
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data["heading"]
        body = slide.placeholders[1].text_frame
        for i, bullet in enumerate(slide_data.get("bullets", [])):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = bullet
    prs.save(path)
    return path


def write_xlsx(path: str, sheet_name: str, headers: list, rows: list) -> str:
    _require(Workbook, "openpyxl", ".xlsx")
    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name
    ws.append(headers)
    for row in rows:
        ws.append(row)
    wb.save(path)
    return path
